#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
from time import time
import os


def gerar_pptx(caminho, lista_arquivos=None):
    # Cria uma apresentação vazia
    prs = Presentation()

    if lista_arquivos:
        arquivos_texto = lista_arquivos
    else:
        # Obtém a lista de arquivos no diretório
        arquivos_texto = [f for f in os.listdir(caminho) if f.endswith('.txt')]

    if arquivos_texto:
        # Itera sobre cada arquivo de texto fornecido
        for idx, arquivo in enumerate(arquivos_texto):
            if lista_arquivos:
                arquivo_full = arquivo
            else:
                arquivo_full = os.path.join(caminho, arquivo)
            try:
                with open(arquivo_full, 'r', encoding='utf-8') as f:
                    # Lê todas as linhas do arquivo
                    lines = f.readlines()

                # Adiciona um slide em branco para separar os arquivos
                if idx > 0:  # Adiciona slide em branco entre os arquivos
                    # Layout em branco (sem título nem conteúdo)
                    slide_empty = prs.slides.add_slide(prs.slide_layouts[6])
                    slide_empty_back = slide_empty.background.fill
                    slide_empty_back.solid()
                    slide_empty_back.fore_color.theme_color = MSO_THEME_COLOR.DARK_1
                # Para dividir os blocos de texto
                bloco = []

                for line in lines:
                    stripped_line = line.strip()

                    # Se a linha não for vazia, adicione ao bloco
                    if stripped_line:
                        bloco.append(stripped_line)
                    else:  # Linha vazia, cria um novo slide com o bloco atual
                        if bloco:  # Se o bloco não estiver vazio
                            # Cria um novo slide com título e conteúdo
                            # Layout de título e conteúdo
                            slide = prs.slides.add_slide(prs.slide_layouts[6])
                            slide_back = slide.background.fill
                            slide_back.solid()
                            slide_back.fore_color.theme_color = MSO_THEME_COLOR.DARK_1

                            content = slide.shapes.add_textbox(
                                left=0, top=0,
                                width=prs.slide_width,
                                height=prs.slide_height
                            )
                            text_frame = content.text_frame
                            text_frame.text = "\n".join(bloco)
                            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                            text_frame.word_wrap = True

                            for ppt in text_frame.paragraphs:
                                ppt.alignment = PP_ALIGN.CENTER
                                ppt.font.size = Pt(50)
                                ppt.font.bold = True
                                ppt.font.color.theme_color = MSO_THEME_COLOR.LIGHT_1

                            # Limpa o bloco para adicionar o próximo
                            bloco = []

                # Bloco de texto não processado, cria um slide
                if bloco:
                    # Layout de título e conteúdo
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    slide_back = slide.background.fill
                    slide_back.solid()
                    slide_back.fore_color.theme_color = MSO_THEME_COLOR.DARK_1

                    content = slide.shapes.add_textbox(
                        left=0, top=0,
                        width=prs.slide_width,
                        height=prs.slide_height
                    )
                    text_frame = content.text_frame
                    text_frame.text = "\n".join(bloco)
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    text_frame.word_wrap = True

                    for ppt in text_frame.paragraphs:
                        ppt.alignment = PP_ALIGN.CENTER
                        ppt.font.size = Pt(50)
                        ppt.font.bold = True
                        ppt.font.color.theme_color = MSO_THEME_COLOR.LIGHT_1

            except Exception as e:
                return "Erro ao processar o arquivo {arquivo}: {e}"

        # Salva a apresentação gerada
        filename = f'{caminho}/apresentacao-{time()}.pptx'
        prs.save(filename)
        return f"Apresentação salva em:\n", filename

    else:
        return "Selecione os arquivos!", None


if __name__ == "__main__":
    # Chame a função passando uma lista de arquivos de texto
    home = os.path.expanduser('~')
    _path_ = f'{home}/Músicas/LETRAS_PROJETOR'
    gerar_pptx(_path_, None)
